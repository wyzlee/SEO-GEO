"""
Génération des présentations SEO-GEO (dark + light)
à partir des templates Wyzlee.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy, os

DARK_TEMPLATE  = "/Users/olivier/Developer/wyz-hub/Wyzlee-Template-Dark.pptx"
LIGHT_TEMPLATE = "/Users/olivier/Developer/wyz-hub/Wyzlee-Template-Light.pptx"
OUT_DIR        = "/Users/olivier/Developer/Chloe/SEO-GEO/guide"

# ── Layout indices ────────────────────────────────────────────────────────────
L_COUVERTURE   = 0   # Eyebrow(10) | Titre(0-CENTER_TITLE) | Sous-titre(0-SUBTITLE)
L_TITRE_CONT   = 1   # Eyebrow(10) | Titre(0) | Contenu(1)
L_SECTION      = 2   # Numero(10)  | Titre(0) | Description(1)
L_DEUX_COL     = 3   # Titre(0)    | ContenuGauche(1) | ContenuDroite(2)
L_COMPARAISON  = 4   # Titre(0)    | ContenuAvant(1) | ContenuApres(2)
L_TITRE_SEUL   = 5   # Eyebrow(10) | Titre(0) | Sous-titre(1)
L_VIERGE       = 6
L_TEXTE_VIS    = 7   # Titre(0)    | Texte(1) | Visuel(2)
L_GRILLE       = 8   # Titre(0)    + shapes CardTitle/CardDesc
L_CITATION     = 9   # Citation(0) | Auteur(1) | Contexte(2)
L_CLOTURE      = 10  # Titre(0-CENTER_TITLE) | Message(0-SUBTITLE) | Contact(1)

# ── Helpers ───────────────────────────────────────────────────────────────────

def _tf_set(tf, text, size=None, bold=None, color_hex=None, align=None, word_wrap=True):
    """Remplace tout le texte d'un text frame."""
    tf.clear()
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    if bold is not None:
        run.font.bold = bold
    if size:
        run.font.size = Pt(size)
    if color_hex:
        run.font.color.rgb = RGBColor.from_string(color_hex)


def _tf_bullets(tf, items, size=11, color_hex=None, word_wrap=True):
    """
    items : list of str  OR  list of (bold_str, rest_str)
    """
    tf.clear()
    tf.word_wrap = word_wrap
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        if isinstance(item, tuple):
            bold_part, rest_part = item
            r1 = p.add_run()
            r1.text = bold_part
            r1.font.bold = True
            r1.font.size = Pt(size)
            if color_hex:
                r1.font.color.rgb = RGBColor.from_string(color_hex)
            if rest_part:
                r2 = p.add_run()
                r2.text = " " + rest_part
                r2.font.bold = False
                r2.font.size = Pt(size)
                if color_hex:
                    r2.font.color.rgb = RGBColor.from_string(color_hex)
        else:
            run = p.add_run()
            run.text = item
            run.font.bold = False
            run.font.size = Pt(size)
            if color_hex:
                run.font.color.rgb = RGBColor.from_string(color_hex)


def _add_textbox(slide, left_in, top_in, w_in, h_in, text,
                 size=10, bold=False, color_hex="8AA0B8",
                 align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in)
    )
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color_hex)
    return txb


def _ph_by_idx(slide, idx):
    """Retourne le placeholder par son idx ou None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_couverture(prs, dark):
    layout = prs.slide_layouts[L_COUVERTURE]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"

    # Eyebrow (idx=10)
    ph = _ph_by_idx(slide, 10)
    if ph:
        _tf_set(ph.text_frame, "WYZLEE — PRODUIT SEO-GEO — AVRIL 2026",
                size=11, bold=True, color_hex=col)

    # Titre (CENTER_TITLE idx=0 → premier placeholder de ce type)
    # Dans la couverture il y a 2 placeholders idx=0 (CENTER_TITLE et SUBTITLE)
    # On les distingue par leur position
    ct = [p for p in slide.placeholders if p.placeholder_format.idx == 0]
    if len(ct) >= 1:
        _tf_set(ct[0].text_frame, "SEO-GEO\nAudit",
                size=48, bold=True, color_hex=txt)
    if len(ct) >= 2:
        _tf_set(ct[1].text_frame,
                "L'audit dual-signal — Google et les moteurs IA\n"
                "(ChatGPT, Perplexity, Claude, Gemini, Copilot)",
                size=16, bold=False, color_hex=muted)

    # Stats sur une ligne en bas
    stats = [
        ("11 phases", "d'analyse"),
        ("Score /100", "actionnable"),
        ("< 10 min", "par audit"),
        ("White-label", "natif FR"),
    ]
    x_start = 1.20
    for i, (val, lbl) in enumerate(stats):
        x = x_start + i * 2.8
        _add_textbox(slide, x, 5.0, 2.5, 0.45, val,
                     size=18, bold=True, color_hex=txt)
        _add_textbox(slide, x, 5.48, 2.5, 0.30, lbl,
                     size=9, bold=False, color_hex=muted)

    return slide


def slide_section(prs, numero, titre, description, dark):
    layout = prs.slide_layouts[L_SECTION]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"

    ph10 = _ph_by_idx(slide, 10)
    if ph10:
        _tf_set(ph10.text_frame, numero, size=56, bold=True, color_hex=col)

    ph0 = _ph_by_idx(slide, 0)
    if ph0:
        _tf_set(ph0.text_frame, titre, size=44, bold=True, color_hex=txt)

    ph1 = _ph_by_idx(slide, 1)
    if ph1:
        _tf_set(ph1.text_frame, description, size=14, bold=False, color_hex=muted)

    return slide


def slide_titre_contenu(prs, eyebrow, titre, items, dark, item_size=11):
    layout = prs.slide_layouts[L_TITRE_CONT]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"

    ph10 = _ph_by_idx(slide, 10)
    if ph10:
        _tf_set(ph10.text_frame, eyebrow, size=10, bold=True, color_hex=col)

    ph0 = _ph_by_idx(slide, 0)
    if ph0:
        _tf_set(ph0.text_frame, titre, size=32, bold=True, color_hex=txt)

    ph1 = _ph_by_idx(slide, 1)
    if ph1:
        _tf_bullets(ph1.text_frame, items, size=item_size, color_hex=muted)

    return slide


def slide_comparaison(prs, titre, label_avant, items_avant,
                                   label_apres, items_apres, dark, item_size=11):
    layout = prs.slide_layouts[L_COMPARAISON]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"
    accent = "EF4444" if dark else "DC2626"
    green  = "39D353" if dark else "16A34A"

    ph0 = _ph_by_idx(slide, 0)
    if ph0:
        _tf_set(ph0.text_frame, titre, size=28, bold=True, color_hex=txt)

    ph1 = _ph_by_idx(slide, 1)
    if ph1:
        _tf_bullets(ph1.text_frame, items_avant, size=item_size, color_hex=muted)

    ph2 = _ph_by_idx(slide, 2)
    if ph2:
        _tf_bullets(ph2.text_frame, items_apres, size=item_size, color_hex=muted)

    # Override les labels des cartes (layout shapes non-editables → textbox overlay)
    _add_textbox(slide, 1.10, 2.00, 4.40, 0.35, label_avant.upper(),
                 size=9, bold=True, color_hex=accent)
    _add_textbox(slide, 7.30, 2.00, 4.40, 0.35, label_apres.upper(),
                 size=9, bold=True, color_hex=green)

    return slide


def slide_deux_colonnes(prs, titre,
                         label_g, items_g,
                         label_d, items_d,
                         dark, item_size=11,
                         col_g_accent=None, col_d_accent=None):
    layout = prs.slide_layouts[L_DEUX_COL]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"

    ph0 = _ph_by_idx(slide, 0)
    if ph0:
        _tf_set(ph0.text_frame, titre, size=28, bold=True, color_hex=txt)

    ph1 = _ph_by_idx(slide, 1)
    if ph1:
        _tf_bullets(ph1.text_frame, items_g, size=item_size, color_hex=muted)

    ph2 = _ph_by_idx(slide, 2)
    if ph2:
        _tf_bullets(ph2.text_frame, items_d, size=item_size, color_hex=muted)

    # Override column labels
    c_g = col_g_accent or col
    c_d = col_d_accent or col
    _add_textbox(slide, 1.10, 2.00, 4.40, 0.35, label_g.upper(),
                 size=9, bold=True, color_hex=c_g)
    _add_textbox(slide, 7.30, 2.00, 4.40, 0.35, label_d.upper(),
                 size=9, bold=True, color_hex=c_d)

    return slide


def slide_grille_cartes(prs, titre, cards, dark):
    """
    cards : list of 6 dicts {title, desc, color_hex}
    """
    layout = prs.slide_layouts[L_GRILLE]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"

    ph0 = _ph_by_idx(slide, 0)
    if ph0:
        _tf_set(ph0.text_frame, titre, size=32, bold=True, color_hex=txt)

    # Positions des 6 cartes (depuis l'inspection du layout)
    positions = [
        (0.80, 1.80), (4.90, 1.80), (9.00, 1.80),
        (0.80, 4.40), (4.90, 4.40), (9.00, 4.40),
    ]
    card_w = 4.00
    card_h = 2.40

    for i, card in enumerate(cards[:6]):
        lx, ly = positions[i]
        c = card.get("color_hex", col)
        t = card.get("title", "")
        d = card.get("desc", "")
        # Titre de carte
        _add_textbox(slide, lx + 0.30, ly + 0.20, card_w - 0.40, 0.35,
                     t, size=10, bold=True, color_hex=c)
        # Description de carte
        _add_textbox(slide, lx + 0.30, ly + 0.60, card_w - 0.40, card_h - 0.75,
                     d, size=9, bold=False, color_hex=muted, word_wrap=True)

    return slide


def slide_citation(prs, citation, auteur, contexte, dark):
    layout = prs.slide_layouts[L_CITATION]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"

    ph0 = _ph_by_idx(slide, 0)
    if ph0:
        _tf_set(ph0.text_frame, citation, size=24, bold=False, color_hex=txt)

    ph1 = _ph_by_idx(slide, 1)
    if ph1:
        _tf_set(ph1.text_frame, auteur, size=14, bold=True, color_hex=col)

    ph2 = _ph_by_idx(slide, 2)
    if ph2:
        _tf_set(ph2.text_frame, contexte, size=12, bold=False, color_hex=muted)

    return slide


def slide_cloture(prs, titre, message, contact, dark):
    layout = prs.slide_layouts[L_CLOTURE]
    slide  = prs.slides.add_slide(layout)
    col    = "6366F1" if dark else "4F46E5"
    txt    = "E2EAF4" if dark else "1A1E2C"
    muted  = "8AA0B8" if dark else "5C6A7E"

    ct = [p for p in slide.placeholders if p.placeholder_format.idx == 0]
    if len(ct) >= 1:
        _tf_set(ct[0].text_frame, titre, size=48, bold=True, color_hex=txt)
    if len(ct) >= 2:
        _tf_set(ct[1].text_frame, message, size=16, bold=False, color_hex=muted)

    ph1 = _ph_by_idx(slide, 1)
    if ph1:
        _tf_set(ph1.text_frame, contact, size=12, bold=False, color_hex=muted)

    # Override CTA button text
    for shape in slide.shapes:
        if shape.name == "CTA" and shape.has_text_frame:
            _tf_set(shape.text_frame, "Lancer un audit → seo-geo-orcin.vercel.app",
                    size=11, bold=True, color_hex="FFFFFF")

    return slide


# ── Contenu des slides ────────────────────────────────────────────────────────

def build_presentation(dark: bool):
    tpl = DARK_TEMPLATE if dark else LIGHT_TEMPLATE
    prs = Presentation(tpl)
    col   = "6366F1" if dark else "4F46E5"
    col2  = "7C3AED"
    txt   = "E2EAF4" if dark else "1A1E2C"
    muted = "8AA0B8" if dark else "5C6A7E"

    # ── 01 — COUVERTURE ───────────────────────────────────────────────────────
    slide_couverture(prs, dark)

    # ── 02 — SECTION : Le contexte ───────────────────────────────────────────
    slide_section(prs,
        "01",
        "Le marché\na basculé",
        "40 % des recherches d'information démarrent sur une IA. "
        "Les outils SEO classiques n'ont pas suivi.",
        dark)

    # ── 03 — TITRE + CONTENU : 4 stats ───────────────────────────────────────
    slide_titre_contenu(prs,
        "CONTEXTE MARCHÉ 2026",
        "4 signaux qui changent la donne",
        [
            ("40 %", "des requêtes d'information démarrent sur une IA, pas sur Google"),
            ("3,6×", "plus de crawls GPTBot vs Googlebot sur les sites analysés"),
            ("76 %", "des citations IA portent sur des contenus de moins de 30 jours"),
            ("3-6 mois", "— demi-vie d'une citation IA, plus courte qu'en SEO classique"),
            "",
            ("Résultat :", "votre client perd simultanément du trafic organique ET de la "
             "visibilité dans les réponses IA. Les deux s'effondrent ensemble."),
        ],
        dark, item_size=12)

    # ── 04 — COMPARAISON : Avant / Après ─────────────────────────────────────
    slide_comparaison(prs,
        "SEO seul ou SEO + GEO ?",
        "AVANT — SEO classique seulement",
        [
            ("Mesure :", "positions Google uniquement"),
            ("Ignore :", "GPTBot, llms.txt, citations IA"),
            ("Aveugle sur :", "ChatGPT Search, Perplexity, Claude"),
            ("Risque :", "trafic organique en baisse sans explication"),
            ("Outils :", "Ahrefs, Semrush — conçus avant 2023"),
            "",
            ("Résultat :", "le client perd de la visibilité et ne sait pas pourquoi"),
        ],
        "APRÈS — Audit dual-signal SEO+GEO",
        [
            ("Mesure :", "Google ET les 5 grands moteurs IA"),
            ("Analyse :", "llms.txt, fichier robot IA, signaux de citation"),
            ("Teste :", "prompts réels sur ChatGPT, Perplexity, Claude"),
            ("Identifie :", "pourquoi une IA cite (ou ne cite pas) votre site"),
            ("Rapport :", "actionnable en français, livré en < 10 min"),
            "",
            ("Résultat :", "le client sait exactement quoi corriger pour remonter"),
        ],
        dark, item_size=11)

    # ── 05 — SECTION : L'Audit ────────────────────────────────────────────────
    slide_section(prs,
        "02",
        "L'Audit en\n11 dimensions",
        "Score sur 100 pts · Rapport FR jargon-free · Moins de 10 minutes",
        dark)

    # ── 06 — TITRE + CONTENU : 11 dimensions ─────────────────────────────────
    slide_titre_contenu(prs,
        "LES 11 DIMENSIONS D'ANALYSE",
        "Ce qu'on mesure sur chaque site",
        [
            ("01  Fondations techniques", "12 pts — URLs, HTTPS, robots.txt, sitemap"),
            ("02  Données structurées",   "15 pts — Schema.org, FAQ, Articles, Produits"),
            ("03  Visibilité IA (GEO) ★", "18 pts — Tests ChatGPT/Perplexity, llms.txt"),
            ("04  Entités & notoriété",   "10 pts — Knowledge Graph, Wikidata, mentions"),
            ("05  E-E-A-T",               "10 pts — Expertise, Autorité, Confiance"),
            ("06  Fraîcheur",             "8 pts  — 76 % citations IA = contenus < 30j"),
            ("07  International",         "8 pts  — Hreflang, multilingue"),
            ("08  Performance web",       "8 pts  — Core Web Vitals, INP, LCP"),
            ("09  Couverture thématique", "6 pts  — Pillar pages, clusters"),
            ("10  Erreurs communes",      "5 pts  — Cannibalisation, thin content"),
            ("11  Synthèse",              "→      — Plan d'action priorisé bénéfice/effort"),
        ],
        dark, item_size=10)

    # ── 07 — DEUX COLONNES : Focus GEO ───────────────────────────────────────
    slide_deux_colonnes(prs,
        "La dimension la plus importante — GEO (18 pts)",
        "CE QUE ÇA TESTE",
        [
            ("Tests de prompts réels :", "on demande à ChatGPT, Perplexity, "
             "Claude et Gemini de répondre aux requêtes de votre domaine. "
             "Apparaissez-vous dans les réponses ?"),
            "",
            ("Fichier llms.txt :", "l'équivalent du robots.txt pour les IA. "
             "Présent ? Correctement configuré ?"),
            "",
            ("Analyse des citations :", "votre marque est-elle citée par les "
             "sources que les IA utilisent — Wikipedia, Wikidata, presse ?"),
        ],
        "POURQUOI C'EST PRIORITAIRE",
        [
            ("18 pts sur 100 :", "c'est le poids le plus élevé de l'audit — "
             "plus que les fondations techniques (12) ou la performance (8)."),
            "",
            ("Impact direct :", "un site bien balisé pour les IA obtient "
             "+132 % de taux de citation dans les réponses GEO."),
            "",
            ("Fraîcheur :", "76 % des citations IA portent sur des contenus "
             "publiés dans les 30 derniers jours. La fraîcheur est un signal "
             "de ranking IA aussi important que l'autorité."),
        ],
        dark, item_size=10,
        col_g_accent="7C3AED", col_d_accent="7C3AED")

    # ── 08 — DEUX COLONNES : Le rapport ───────────────────────────────────────
    slide_deux_colonnes(prs,
        "Le rapport livré au client",
        "LE SCORE & LE PLAN D'ACTION",
        [
            ("Score sur 100 :", "un chiffre unique, clair, comparable d'un "
             "audit à l'autre. Le client comprend immédiatement où il en est."),
            "",
            ("Détail par dimension :", "score de chaque phase avec les "
             "findings les plus impactants."),
            "",
            ("Synthèse — 11ème dimension :", "les 5 recommandations "
             "les plus impactantes, classées par rapport bénéfice / effort. "
             "Ce que le client doit faire en premier."),
            "",
            ("Stats sourcées :", "chaque chiffre cité est traçable "
             "à une source consultée et datée."),
        ],
        "LE FORMAT DE LIVRAISON",
        [
            ("Lien web partageable :", "URL tokenisée, durée configurable. "
             "Le client clique, il voit le rapport depuis n'importe quel device."),
            "",
            ("Export PDF :", "qualité impression, idéal pour réunions "
             "client et archivages contractuels."),
            "",
            ("White-label natif :", "votre logo (ou celui de votre client). "
             "Zéro mention SEO-GEO dans le document livré."),
            "",
            ("Français, sans jargon :", "rédigé pour être compris par un "
             "décideur non-technique. Pas de termes SEO obscurs."),
        ],
        dark, item_size=10)

    # ── 09 — COMPARAISON : URL vs Code ───────────────────────────────────────
    slide_comparaison(prs,
        "Deux façons de lancer un audit",
        "URL LIVE  (V1 — disponible)",
        [
            ("Comment :", "entrez l'URL d'un site en production. "
             "L'audit analyse le site tel que les moteurs le voient."),
            "",
            ("Couvre :", "contenus visibles, performances réelles, "
             "signaux IA en live, robots.txt, sitemap.xml"),
            "",
            ("Idéal pour :", "audits clients sur sites existants, "
             "suivis trimestriels, prospection commerciale"),
            "",
            ("→", "Arrivez en RDV avec un rapport sous le nom du prospect."),
        ],
        "CODE SOURCE  (V2 — prochainement)",
        [
            ("Comment :", "uploadez un ZIP ou connectez un repo GitHub. "
             "L'audit analyse le code avant mise en ligne."),
            "",
            ("Couvre :", "structure des composants, balisage statique, "
             "schémas JSON-LD, configuration SEO dans le code"),
            "",
            ("Idéal pour :", "studios dev qui livrent un site optimisé "
             "dès le départ, audits de refonte, pre-launch"),
            "",
            ("→", "Livrez un asset mieux foutu à votre client — "
             "avec le rapport d'audit en preuve."),
        ],
        dark, item_size=10)

    # ── 10 — SECTION : Pour qui & Combien ────────────────────────────────────
    slide_section(prs,
        "03",
        "Pour qui\net combien ?",
        "3 personas · 6 packages · De 1 500 € à 15 000 €/mois",
        dark)

    # ── 11 — GRILLE CARTES : Personas + Offres ───────────────────────────────
    v1c = "6366F1" if dark else "4F46E5"
    v2c = "7C3AED"
    v3c = "3882F6" if dark else "2563EB"
    v4c = "8B5CF6"
    v5c = "F59E0B" if dark else "D97706"
    v6c = "EF4444"

    slide_grille_cartes(prs,
        "Pour qui et combien ?",
        [
            {"title": "Agence SEO", "color_hex": v1c,
             "desc": "5-30 pers. Veut ajouter une offre GEO sans recruter. "
                     "Offre : SEO+GEO Add-on (+25 %) ou White-label (-40 %)."},
            {"title": "Directeur marketing B2B SaaS", "color_hex": v2c,
             "desc": "Trafic organique qui stagne. Veut comprendre pourquoi "
                     "ses concurrents apparaissent dans Perplexity. "
                     "Offre : Tripwire 1 500-3 500 € → Retainer."},
            {"title": "Studio dev / Freelance", "color_hex": v3c,
             "desc": "Livre des sites clients. Veut auditer avant release "
                     "(V2 code source). "
                     "Offre : Tripwire par projet ou White-label."},
            {"title": "Tripwire Audit", "color_hex": v4c,
             "desc": "1 500–3 500 € one-shot. Audit 11 dimensions + debrief 1h. "
                     "Livraison 3-5 jours. Taux de conversion cible > 35 %."},
            {"title": "Retainer Starter/Growth/Enterprise", "color_hex": v5c,
             "desc": "2 500–15 000 €/mois. Mensuel ou trimestriel. "
                     "Starter 3 mois / Growth 6 mois / Enterprise 12 mois. "
                     "Tracking continu inclus en Growth+."},
            {"title": "White-label & Add-on", "color_hex": v6c,
             "desc": "White-label : wholesale -40 à -60 %. "
                     "SEO+GEO Add-on : +25 % sur retainer SEO existant. "
                     "Exécution sous la marque de l'agence partenaire."},
        ],
        dark)

    # ── 12 — DEUX COLONNES : Roadmap ─────────────────────────────────────────
    slide_deux_colonnes(prs,
        "Agency Tool V1 → Self-serve SaaS V2",
        "MAINTENANT — V1 AGENCY MODE",
        [
            ("Dashboard interne :", "Olivier lance et suit les audits. "
             "Le client reçoit uniquement le rapport."),
            ("Audit URL live :", "11 dimensions, score 100 pts, < 10 min."),
            ("Rapport white-label :", "web partageable + PDF, en français."),
            ("Pricing à la prestation :", "tripwire + retainer."),
            "",
            "Objectif 90 jours :",
            ("→", "10+ audits livrés à de vrais clients"),
            ("→", "3+ conversions tripwire → retainer"),
            ("→", "CSAT ≥ 4.3/5 sur les rapports"),
        ],
        "PROCHAINEMENT — V2 SELF-SERVE",
        [
            ("Signup public :", "inscription directe, onboarding automatisé."),
            ("Dashboard client :", "historique audits, comparaison, KPI continu."),
            ("Plans Stripe :", "Free (1/mois) · Pro (10/mois) · Agence (illimité)."),
            ("Code source :", "upload ZIP + GitHub connect, audit pre-launch."),
            ("Alertes :", "notification automatique si score régressé."),
            "",
            "Métriques V2 :",
            ("→", "≥ 50 signups/mois"),
            ("→", "≥ 10 % conversion Free → Pro"),
        ],
        dark, item_size=10,
        col_g_accent=col, col_d_accent=col2)

    # ── 13 — CLÔTURE ─────────────────────────────────────────────────────────
    slide_cloture(prs,
        "Prêt à lancer\nle premier audit ?",
        "Entrez une URL. Obtenez un rapport complet en moins de 10 minutes.\n"
        "White-label, en français, livrable client immédiatement.",
        "wyzlee.com / olivier.podio@gmail.com",
        dark)

    return prs


# ── Main ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    os.makedirs(OUT_DIR, exist_ok=True)

    print("Génération de la version DARK…")
    prs_dark = build_presentation(dark=True)
    out_dark = os.path.join(OUT_DIR, "SEO-GEO-Presentation-Dark.pptx")
    prs_dark.save(out_dark)
    print(f"  ✓ {out_dark}  ({len(prs_dark.slides)} slides)")

    print("Génération de la version LIGHT…")
    prs_light = build_presentation(dark=False)
    out_light = os.path.join(OUT_DIR, "SEO-GEO-Presentation-Light.pptx")
    prs_light.save(out_light)
    print(f"  ✓ {out_light}  ({len(prs_light.slides)} slides)")

    print("\nDone.")
